import os
from pptx import Presentation
from openpyxl import Workbook

def pptx_to_excel(pptx_path, output_xlsx):
    """
    将PPTX文件转换为Excel文件，每页PPT内容单独存储在一个Excel sheet中
    
    参数:
        pptx_path: 输入的PPTX文件路径
        output_xlsx: 输出的Excel文件路径
    """
    # 创建Excel工作簿
    wb = Workbook()
    # 删除默认创建的Sheet
    wb.remove(wb.active)
    
    # 加载PPTX文件
    prs = Presentation(pptx_path)
    
    # 遍历每页PPT
    for i, slide in enumerate(prs.slides):
        # 为每页PPT创建一个新的Sheet
        sheet = wb.create_sheet(title=f"Slide_{i+1}")
        
        # 提取PPT页面中的内容
        row = 1
        for shape in slide.shapes:
            # 处理文本内容
            if hasattr(shape, "text"):
                sheet.cell(row=row, column=1, value=shape.text)
                row += 1
            
            # 处理表格内容
            if shape.has_table:
                table = shape.table
                for r_idx, row_data in enumerate(table.rows):
                    for c_idx, cell in enumerate(row_data.cells):
                        if cell.text:
                            sheet.cell(row=row+r_idx, column=c_idx+1, value=cell.text)
                row += len(table.rows)
    
    # 处理单位转换
    for sheet in wb:
        try:
            # 通过表头关键字查找目标单元格
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        # 检查尺寸相关表头
                        if 'Product Dimensions (mm*mm*mm)' in cell.value:
                            cell.value = 'Product Dimensions (inch)'
                            target_cell = sheet.cell(row=cell.row, column=cell.column+1)
                            if target_cell.value:
                                try:
                                    # 转换长*宽*高尺寸
                                    dimensions = [round(float(d)/25.4, 1) for d in str(target_cell.value).split('*')]
                                    target_cell.value = '*'.join(map(str, dimensions))
                                except (ValueError, AttributeError):
                                    pass
                        elif 'Active Display Area (mm*mm)' in cell.value:
                            cell.value = 'Active Display Area (inch)'
                            target_cell = sheet.cell(row=cell.row, column=cell.column+1)
                            if target_cell.value:
                                try:
                                    # 转换长*宽尺寸
                                    dimensions = [round(float(d)/25.4, 1) for d in str(target_cell.value).split('*')]
                                    target_cell.value = '*'.join(map(str, dimensions))
                                except (ValueError, AttributeError):
                                    pass
                        
                        # 检查温度相关表头
                        elif 'Working Temperature (˚C)' in cell.value or 'Working Temperature (℃)' in cell.value:
                            cell.value = 'Working Temperature (℉)'
                            target_cell = sheet.cell(row=cell.row, column=cell.column+1)
                            if target_cell.value:
                                try:
                                    temp_str = str(target_cell.value).strip()
                                    # 处理温度范围格式(支持'0~40'或'0-40'格式)
                                    if '~' in temp_str or '-' in temp_str:
                                        separator = '~' if '~' in temp_str else '-'
                                        temps = [t.strip() for t in temp_str.split(separator)]
                                        f_temps = [str(int(float(t)*9/5 + 32)) for t in temps]
                                        target_cell.value = f' {separator} '.join(f_temps)
                                    # 处理单个温度值
                                    else:
                                        f_temp = str(int(float(temp_str)*9/5 + 32))
                                        target_cell.value = f_temp
                                except (ValueError, AttributeError) as e:
                                    print(f"温度转换错误: {e}")
        except Exception as e:
            print(f"处理工作表 {sheet.title} 时出错: {e}")
            continue
    
    # 保存Excel文件
    try:
        wb.save(output_xlsx)
        print(f"成功将PPTX文件转换为Excel文件: {output_xlsx}")
    except PermissionError:
        print(f"错误: 没有权限写入到目录 {os.path.dirname(output_xlsx)}，请检查目录权限或尝试其他位置")
        sys.exit(1)
    except Exception as e:
        print(f"保存Excel文件时出错: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) != 3:
        print("使用方法: python ppt2excel.py <输入PPTX文件路径> <输出Excel文件路径>")
        sys.exit(1)
    
    input_pptx = sys.argv[1]
    output_xlsx = sys.argv[2]
    
    if not os.path.exists(input_pptx):
        print(f"错误: 输入的PPTX文件不存在 - {input_pptx}")
        sys.exit(1)
    
    pptx_to_excel(input_pptx, output_xlsx)