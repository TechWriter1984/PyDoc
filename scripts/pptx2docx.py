from pptx import Presentation

def pptx_to_text(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

text_content = pptx_to_text(r"D:\对外文档归档\AP\五代基站\原文档\diagrams.pptx")
print(text_content)