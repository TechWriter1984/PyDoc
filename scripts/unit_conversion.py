from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re
import os

MM_PER_INCH = 25.4

# ========= Utility =========

def mm_to_inch(mm_val: float) -> float:
    return mm_val / MM_PER_INCH

def format_inch(val: float) -> str:
    s = f"{val:.2f}"
    if s.endswith("00"):
        return str(int(round(val)))
    return s

def convert_mm_pattern(text: str) -> str:
    parts = re.split(r"[*×]", text.strip())
    if not (2 <= len(parts) <= 3):
        return text

    try:
        new_vals = [format_inch(mm_to_inch(float(p.strip()))) for p in parts]
    except:
        return text

    sep = "*" if "*" in text else "×"
    return sep.join(new_vals)

def convert_celsius_range(text: str) -> str:
    m = re.match(r"\s*([-+]?\d+(\.\d+)?)\s*[~-]\s*([-+]?\d+(\.\d+)?)\s*$", text)
    if not m:
        return text

    c1, c2 = float(m.group(1)), float(m.group(3))

    def c2f(c):
        return c * 9/5 + 32

    f1, f2 = round(c2f(c1)), round(c2f(c2))
    sep = "~" if "~" in text else "-"
    return f"{f1}{sep}{f2}"

def set_cell_font(cell, font="Montserrat-Light", size=7):
    if not cell.text_frame:
        return
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)

def set_shape_font(shape, font="Montserrat-Light", size=7):
    if not shape.has_text_frame:
        return
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)

# ========= NEW: recursive shape iterator =========

def iter_all_shapes(shapes):
    """
    深度遍历所有 shape，包括 GroupShape 内部的 shapes
    """
    for shape in shapes:
        yield shape

        # 组合形状 (GroupShape)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP = 6
            for sub in iter_all_shapes(shape.shapes):
                yield sub



# ========= Main Logic =========

def process_pptx(input_path: str):
    prs = Presentation(input_path)

    HEAD_DIM = "Product Dimensions (mm*mm*mm)"
    HEAD_AREA = "Active Display Area (mm*mm)"
    HEAD_TEMP_C = {
        "Working Temperature (˚C)",
        "Working Temperature (°C)",
        "Working Temperature (C)",
        "Working Temperature (℃)",
    }

    for slide in prs.slides:

        # ========== Process tables ==========
        for shape in iter_all_shapes(slide.shapes):
            if not shape.has_table:
                continue

            table = shape.table

            for row in table.rows:
                if len(row.cells) < 2:
                    continue

                hcell, vcell = row.cells[0], row.cells[1]
                header = hcell.text.strip()
                value = vcell.text.strip()

                if header == HEAD_DIM:
                    new_val = convert_mm_pattern(value)
                    vcell.text = new_val
                    hcell.text = "Product Dimensions (inch*inch*inch)"
                    set_cell_font(hcell)
                    set_cell_font(vcell)

                elif header == HEAD_AREA:
                    new_val = convert_mm_pattern(value)
                    vcell.text = new_val
                    hcell.text = "Active Display Area (inch*inch)"
                    set_cell_font(hcell)
                    set_cell_font(vcell)

                elif header in HEAD_TEMP_C:
                    new_val = convert_celsius_range(value)
                    vcell.text = new_val
                    hcell.text = "Working Temperature (°F)"
                    set_cell_font(hcell)
                    set_cell_font(vcell)

        # ========== Process standalone text shapes (mm values) ==========
        for shape in iter_all_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()

            if re.fullmatch(r"\d+(\.\d+)?", text):
                mm_val = float(text)
                inch_text = format_inch(mm_to_inch(mm_val))
                shape.text = inch_text
                set_shape_font(shape)

    # ===== Save new file =====
    base, ext = os.path.splitext(input_path)
    output_path = base + "-US" + ext
    prs.save(output_path)

    print(f"Converted successfully → {output_path}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python unit_conversion.py input.pptx")
        sys.exit(1)

    process_pptx(sys.argv[1])
